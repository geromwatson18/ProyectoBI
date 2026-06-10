# -*- coding: utf-8 -*-
"""
Genera la presentacion .pptx de "Las Salsas de Lucho" a partir de las graficas
ya producidas (evidencia/capturas/*.png y docs/diagramas/*.png).

Uso:  python generar_presentacion.py
Requiere: python-pptx  (pip install python-pptx)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
CAPT = os.path.normpath(os.path.join(HERE, "..", "evidencia", "capturas"))
DIAG = os.path.normpath(os.path.join(HERE, "..", "docs", "diagramas"))
OUT  = os.path.join(HERE, "Presentacion_Las_Salsas_de_Lucho.pptx")

# Paleta (tema salsa: rojos)
ROJO   = RGBColor(0xC0, 0x39, 0x2B)
VINO   = RGBColor(0x7B, 0x24, 0x1C)
CREMA  = RGBColor(0xFC, 0xF6, 0xF5)
ORO    = RGBColor(0xE6, 0x7E, 0x22)
TINTA  = RGBColor(0x2B, 0x2B, 0x2B)
GRIS   = RGBColor(0x70, 0x70, 0x70)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
HEAD = "Cambria"; BODY = "Calibri"

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def box(s, x, y, w, h):
    return s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def text(tb, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    first = True
    for r in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space)
        if "bullet" in r and r["bullet"]:
            p.text = "•  " + r["t"]
        else:
            p.text = r["t"]
        f = p.runs[0].font
        f.size = Pt(r.get("sz", 16)); f.bold = r.get("b", False)
        f.italic = r.get("i", False); f.name = r.get("font", BODY)
        f.color.rgb = r.get("c", TINTA)
    return tb

def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def title_bar(s, t, sub=None):
    text(box(s, 0.6, 0.35, 12.1, 1.0),
         [{"t": t, "sz": 30, "b": True, "font": HEAD, "c": VINO}])
    if sub:
        text(box(s, 0.62, 1.15, 12.1, 0.5),
             [{"t": sub, "sz": 15, "i": True, "c": GRIS}])

def pic_fit(s, path, x, y, w, h):
    """Inserta imagen ajustada dentro de la caja (x,y,w,h) en pulgadas, centrada."""
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw/ih; boxar = w/h
    if ar > boxar:
        nw = w; nh = w/ar
    else:
        nh = h; nw = h*ar
    nx = x + (w-nw)/2; ny = y + (h-nh)/2
    s.shapes.add_picture(path, Inches(nx), Inches(ny), Inches(nw), Inches(nh))

def stat(s, x, y, w, big, label, color=ROJO):
    text(box(s, x, y, w, 0.9), [{"t": big, "sz": 40, "b": True, "c": color, "font": HEAD}], align=PP_ALIGN.CENTER)
    text(box(s, x, y+0.85, w, 0.7), [{"t": label, "sz": 13, "c": TINTA}], align=PP_ALIGN.CENTER)

# ============================ 1 · PORTADA ============================
s = slide(); bg(s, VINO)
rect(s, 0, 0, 13.333, 0.001, VINO)  # no-op para mantener fondo
text(box(s, 1.0, 2.1, 11.3, 2.0), [
    {"t": "🌶️ Las Salsas de Lucho", "sz": 54, "b": True, "font": HEAD, "c": BLANCO},
    {"t": "Solución integral de Inteligencia de Negocios", "sz": 26, "c": CREMA},
])
text(box(s, 1.02, 4.6, 11.3, 1.8), [
    {"t": "Proyecto 02 · Curso de Inteligencia de Negocios", "sz": 16, "c": CREMA},
    {"t": "Tecnológico de Costa Rica — ATI", "sz": 16, "c": CREMA},
    {"t": "Gerom Watson Araya · carné 2017168284 · 2026", "sz": 14, "i": True, "c": ORO},
])

# ============================ 2 · CONTEXTO ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "La organización", "PYME comercial de salsas artesanales · Cartago, Costa Rica")
text(box(s, 0.7, 1.9, 6.4, 4.8), [
    {"t": "Microempresa de alimentos fundada en 2021.", "sz": 17, "bullet": True},
    {"t": "Produce salsas y aderezos artesanales en botella de 200 ml.", "sz": 17, "bullet": True},
    {"t": "Precio de venta: ₡3.000 por botella.", "sz": 17, "bullet": True},
    {"t": "Canal principal: mayoreo a ~80 puntos de venta (pulperías, mini súper, sodas, restaurantes) en todo el país.", "sz": 17, "bullet": True},
    {"t": "Otros canales: tienda física, ferias del agricultor y pedidos en línea.", "sz": 17, "bullet": True},
    {"t": "Registro en facturación sencilla + hojas de cálculo → problemas de calidad de datos.", "sz": 17, "bullet": True},
], space=10)
rect(s, 7.5, 2.1, 5.2, 4.2, CREMA)
stat(s, 7.6, 2.5, 2.45, "₡3.000", "precio · 200 ml")
stat(s, 10.15, 2.5, 2.45, "~80", "puntos de venta/mes", color=ORO)
stat(s, 7.6, 4.4, 2.45, "4", "canales de venta")
stat(s, 10.15, 4.4, 2.45, "2021", "año de fundación", color=ORO)

# ============================ 3 · PROBLEMA Y PREGUNTAS ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "El problema y las preguntas de negocio")
text(box(s, 0.7, 1.8, 12.0, 1.1), [
    {"t": "El propietario decide por intuición: no sabe qué productos y canales dejan utilidad, cuánto inventario lento acumula, ni cómo planificar la producción.", "sz": 16, "i": True, "c": TINTA}])
preguntas = [
    ("PN1", "¿Qué productos concentran ventas y margen, y cuáles rotan poco?"),
    ("PN2", "¿Qué canal es más rentable y cómo evoluciona?"),
    ("PN3", "¿Cómo es la estacionalidad de la demanda?"),
    ("PN4", "¿Las promociones ayudan o erosionan el margen?"),
    ("PN5", "(Exploratoria) ¿Hay dependencia de pocos clientes o de un canal? Geografía."),
]
y = 3.05
for code, q in preguntas:
    rect(s, 0.7, y, 1.0, 0.62, ROJO)
    text(box(s, 0.7, y, 1.0, 0.62), [{"t": code, "sz": 16, "b": True, "c": BLANCO}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(box(s, 1.85, y, 10.8, 0.62), [{"t": q, "sz": 16, "c": TINTA}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.78

# ============================ 4 · ARQUITECTURA ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "Arquitectura de la solución", "Trazabilidad: fuente → ETL → modelo dimensional → tablero")
pasos = [
    ("Fuente operacional", "CSV / PostgreSQL\n(salsas_op) · 9 entidades", ROJO),
    ("ETL en Python", "pandas + SQLAlchemy\n10 reglas + calidad", ORO),
    ("Modelo dimensional", "Estrella (salsas_dw)\n2 hechos · 6 dimensiones", VINO),
    ("Tablero analítico", "Streamlit + Plotly\n6 vistas", ROJO),
]
x = 0.7
for i, (t, d, col) in enumerate(pasos):
    rect(s, x, 2.7, 2.75, 2.0, col)
    text(box(s, x+0.12, 2.85, 2.5, 1.85),
         [{"t": t, "sz": 17, "b": True, "c": BLANCO, "font": HEAD},
          {"t": d, "sz": 13, "c": CREMA}], anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        text(box(s, x+2.78, 3.2, 0.5, 1.0), [{"t": "➜", "sz": 30, "b": True, "c": GRIS}], align=PP_ALIGN.CENTER)
    x += 3.13
text(box(s, 0.7, 5.2, 12.0, 1.2), [
    {"t": "Herramienta de ETL distinta de Power BI / Tableau (requisito del curso).", "sz": 15, "bullet": True},
    {"t": "Pipeline reproducible y agnóstico del motor (PostgreSQL en producción, SQLite para verificación).", "sz": 15, "bullet": True},
])

# ============================ 5 · FUENTE OPERACIONAL ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "Fuente operacional", "9 entidades relacionadas · 2 años de historia · problemas de calidad reales")
filas = [("categorias","3"),("productos","10"),("clientes","137"),("canales","4"),
         ("vendedores","4"),("promociones","5"),("ventas","3.655"),
         ("detalle_ventas","5.055"),("movimientos_inventario","240")]
y=2.0
for i,(e,n) in enumerate(filas):
    yy = 2.0 + (i%5)*0.62
    xx = 0.7 if i<5 else 6.6
    text(box(s, xx, yy, 4.0, 0.55), [{"t": e, "sz": 15, "c": TINTA}], anchor=MSO_ANCHOR.MIDDLE)
    text(box(s, xx+4.0, yy, 1.2, 0.55), [{"t": n, "sz": 15, "b": True, "c": ROJO}], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
pic_fit(s, os.path.join(DIAG, "modelo_operacional.png"), 6.4, 4.7, 6.6, 2.6)
text(box(s, 6.4, 4.45, 6.6, 0.3), [{"t": "Modelo operacional (origen)", "sz": 12, "i": True, "c": GRIS}], align=PP_ALIGN.CENTER)

# ============================ 6 · MODELO DIMENSIONAL ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "Modelo dimensional — esquema estrella", "2 hechos · 6 dimensiones · dim_tiempo · llaves subrogadas · dimensiones conformadas")
pic_fit(s, os.path.join(DIAG, "modelo_dimensional.png"), 0.7, 1.9, 8.2, 5.2)
text(box(s, 9.1, 2.1, 3.7, 5.0), [
    {"t": "fact_ventas", "sz": 16, "b": True, "c": ROJO, "font": HEAD},
    {"t": "grano: línea de factura (4.989 filas)", "sz": 13},
    {"t": "fact_inventario", "sz": 16, "b": True, "c": ROJO, "font": HEAD},
    {"t": "grano: snapshot mensual/producto (240)", "sz": 13},
    {"t": "Conformadas: dim_tiempo y dim_producto", "sz": 13, "bullet": True},
    {"t": "Jerarquías de tiempo, producto y geografía", "sz": 13, "bullet": True},
    {"t": "≥17 medidas documentadas", "sz": 13, "bullet": True},
], space=8)

# ============================ 7 · ETL ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "Proceso ETL — 10 reglas no triviales", "Validación de calidad ANTES del ETL; transformaciones documentadas y trazables")
reglas = [
    "R1 · dim_tiempo (calendario + temporada)",
    "R2 · llaves subrogadas en dimensiones",
    "R3 · homologación de catálogos (23 prov. / 19 canal)",
    "R4 · antigüedad de cliente (imputa 16 nulos)",
    "R5 · segmentación RFM de clientes",
    "R6 · outliers de cantidad por IQR (24)",
    "R7 · deduplicación de líneas (24)",
    "R8 · medidas: neto, costo, margen (imputa 212 costos)",
    "R9 · integridad referencial → 'No identificado' (17)",
    "R10 · tipado y estandarización de formatos",
]
for i,r in enumerate(reglas):
    xx = 0.7 if i<5 else 6.85
    yy = 2.0 + (i%5)*0.7
    rect(s, xx, yy, 5.9, 0.58, CREMA)
    text(box(s, xx+0.15, yy, 5.7, 0.58), [{"t": r, "sz": 13.5, "c": TINTA}], anchor=MSO_ANCHOR.MIDDLE)
text(box(s, 0.7, 5.75, 12.0, 0.8), [
    {"t": "Evidencia: bitácoras en evidencia/etl_logs/ · resultado: fact_ventas = 4.989 filas.", "sz": 14, "i": True, "c": VINO}])

# ============================ 8 · CALIDAD ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "Validación de calidad de datos (pre-ETL)")
pic_fit(s, os.path.join(CAPT, "07_calidad.png"), 0.7, 1.9, 8.0, 5.0)
text(box(s, 8.9, 2.2, 4.0, 4.6), [
    {"t": "Se perfila la calidad ANTES de transformar.", "sz": 15, "bullet": True},
    {"t": "Nulos en fecha, canal y costo.", "sz": 15, "bullet": True},
    {"t": "Cantidades atípicas y negativas.", "sz": 15, "bullet": True},
    {"t": "Duplicados y FK rotas.", "sz": 15, "bullet": True},
    {"t": "Cada problema lo corrige una regla R1–R10.", "sz": 15, "bullet": True, "c": VINO, "b": True},
], space=10)

# ============================ 9 · DASHBOARD ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "Tablero analítico — 6 vistas", "Streamlit + Plotly · segmentadores por año, canal y categoría")
pic_fit(s, os.path.join(CAPT, "01_pareto_productos.png"), 0.7, 1.9, 8.1, 5.0)
vistas = ["1 · Resumen ejecutivo","2 · Productos y rotación","3 · Canales y vendedores",
          "4 · Estacionalidad","5 · Promociones","6 · Clientes y geografía"]
text(box(s, 9.0, 2.2, 4.0, 4.8), [{"t": v, "sz": 15, "bullet": True} for v in vistas], space=12)

# ============================ 10 · PN1 ============================
def slide_pn(titulo, sub, img, lectura):
    s = slide(); bg(s, BLANCO)
    title_bar(s, titulo, sub)
    pic_fit(s, os.path.join(CAPT, img), 0.7, 1.85, 9.0, 5.1)
    rect(s, 9.9, 2.4, 3.0, 3.8, CREMA)
    text(box(s, 10.05, 2.55, 2.75, 3.6), lectura, space=9)
    return s

slide_pn("PN1 · Productos y rotación",
         "Las 3 estrella concentran el 85,2% del ingreso",
         "01_pareto_productos.png",
         [{"t":"Habanero 33% · Chilero 26% · Ajo 26%.","sz":14,"bullet":True},
          {"t":"7 especialidades < 15% del ingreso.","sz":14,"bullet":True},
          {"t":"Baja rotación: ~390 días de inventario.","sz":14,"bullet":True,"c":ROJO,"b":True}])

# ============================ 11 · PN2 ============================
slide_pn("PN2 · Rentabilidad por canal",
         "Mayoreo = 90,9% del ingreso (₡44,4 M)",
         "02_canal.png",
         [{"t":"~80 puntos de venta sostienen el negocio.","sz":14,"bullet":True},
          {"t":"Márgenes parejos: 57–59% entre canales.","sz":14,"bullet":True},
          {"t":"Riesgo: casi todo pasa por un solo canal.","sz":14,"bullet":True,"c":ROJO,"b":True}])

# ============================ 12 · PN3 ============================
slide_pn("PN3 · Estacionalidad",
         "Ingreso promedio diario por mes",
         "03_estacionalidad.png",
         [{"t":"Picos en noviembre–diciembre.","sz":14,"bullet":True},
          {"t":"Piso en enero–febrero.","sz":14,"bullet":True},
          {"t":"Anticipar producción de estrella antes de nov.","sz":14,"bullet":True,"c":VINO,"b":True}])

# ============================ 13 · PN4 ============================
slide_pn("PN4 · Promociones",
         "'Liquidación lenta' destruye margen",
         "04_promociones.png",
         [{"t":"Sin promo: 59,7% de margen.","sz":14,"bullet":True},
          {"t":"Liquidación lenta: 45,6% y solo 0,4% del ingreso.","sz":14,"bullet":True},
          {"t":"Eliminar o sustituir por combos.","sz":14,"bullet":True,"c":ROJO,"b":True}])

# ============================ 14 · PN5 ============================
slide_pn("PN5 · Clientes y geografía",
         "Ingreso bien repartido entre los ~80 puntos",
         "05_clientes.png",
         [{"t":"Top-5 clientes = solo 6,8%: sin dependencia de pocos.","sz":14,"bullet":True},
          {"t":"GAM 65,5% · zonas costeras 34,1%.","sz":14,"bullet":True},
          {"t":"El riesgo es de canal y surtido, no de cliente.","sz":14,"bullet":True,"c":VINO,"b":True}])

# ============================ 15 · HALLAZGOS Y RECOMENDACIONES ============================
s = slide(); bg(s, BLANCO)
title_bar(s, "Hallazgos y recomendaciones accionables")
text(box(s, 0.7, 1.9, 6.0, 5.0), [
    {"t":"Hallazgos","sz":20,"b":True,"c":VINO,"font":HEAD},
    {"t":"85% del ingreso en 3 productos.","sz":15,"bullet":True},
    {"t":"91% del ingreso por mayoreo (~80 puntos).","sz":15,"bullet":True},
    {"t":"Demanda estacional (picos nov–dic).","sz":15,"bullet":True},
    {"t":"Promos que erosionan margen.","sz":15,"bullet":True},
    {"t":"Sin dependencia de pocos clientes.","sz":15,"bullet":True},
], space=9)
text(box(s, 6.9, 1.9, 6.0, 5.0), [
    {"t":"Recomendaciones","sz":20,"b":True,"c":ROJO,"font":HEAD},
    {"t":"Racionalizar el catálogo lento.","sz":15,"bullet":True},
    {"t":"Planificar producción por estacionalidad.","sz":15,"bullet":True},
    {"t":"Revisar política de promociones.","sz":15,"bullet":True},
    {"t":"Reducir dependencia del mayoreo (tienda y línea).","sz":15,"bullet":True},
    {"t":"Cerrar catálogos para mejorar la calidad de datos.","sz":15,"bullet":True},
], space=9)

# ============================ 16 · CIERRE ============================
s = slide(); bg(s, VINO)
text(box(s, 1.0, 2.3, 11.3, 2.6), [
    {"t":"Limitaciones y mejoras futuras","sz":30,"b":True,"font":HEAD,"c":BLANCO},
    {"t":"Datos sintéticos/anonimizados; margen bruto (sin costos indirectos); RFM simplificado.","sz":16,"c":CREMA},
    {"t":"Futuro: integrar costos indirectos, pronóstico de demanda y datos en tiempo real.","sz":16,"c":CREMA},
])
text(box(s, 1.02, 5.2, 11.3, 1.6), [
    {"t":"Trabajo colaborativo y reproducible en GitHub:","sz":16,"b":True,"c":ORO},
    {"t":"github.com/geromwatson18/ProyectoBI","sz":18,"c":BLANCO},
    {"t":"¡Gracias!","sz":20,"b":True,"c":BLANCO,"font":HEAD},
])

prs.save(OUT)
print("Presentacion generada:", OUT, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
